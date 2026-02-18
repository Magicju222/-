"""
PPT文档生成器
生成视觉呈现效果佳、重点突出的PPT演示文稿
"""

import json
from typing import Dict, List, Optional, Any
from io import BytesIO
from datetime import datetime
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from analyzer import AnalysisResult, AnalysisDimension


class PPTDocumentGenerator:
    """
    PPT文档生成器
    生成专业的数据分析汇报PPT
    """
    
    def __init__(self):
        self.prs = None
        self.slide_width = Inches(13.333)
        self.slide_height = Inches(7.5)
    
    def generate(self,
                 analysis_result: AnalysisResult,
                 chart_images: List[str] = None,
                 title: str = "数据分析汇报") -> BytesIO:
        """
        生成PPT演示文稿
        
        Args:
            analysis_result: 分析结果
            chart_images: 图表图像路径列表
            title: 报告标题
        
        Returns:
            PPT文档字节流
        """
        self.prs = Presentation()
        self.prs.slide_width = self.slide_width
        self.prs.slide_height = self.slide_height
        
        # 1. 标题页
        self._add_title_slide(title)
        
        # 2. 执行摘要页
        self._add_summary_slide(analysis_result)
        
        # 3. 数据概览页
        self._add_data_overview_slide(analysis_result)
        
        # 4. 核心发现页（每页一个洞察）
        for i, insight in enumerate(analysis_result.insights[:3]):
            chart_path = chart_images[i] if chart_images and i < len(chart_images) else None
            self._add_insight_slide(insight, i+1, chart_path)
        
        # 5. 结论建议页
        self._add_conclusion_slide(analysis_result)
        
        # 6. 感谢页
        self._add_thank_you_slide()
        
        # 保存到内存
        output = BytesIO()
        self.prs.save(output)
        output.seek(0)
        
        return output
    
    def _add_title_slide(self, title: str):
        """添加标题页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "基于AI的数据分析"
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 日期
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.2), Inches(12.333), Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime('%Y年%m月%d日')
        date_frame.paragraphs[0].font.size = Pt(16)
        date_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_summary_slide(self, result: AnalysisResult):
        """添加执行摘要页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "执行摘要"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # 关键数字（大字突出显示）
        metrics = [
            (str(len(result.merged_dimensions)), "分析维度"),
            (str(len(result.insights)), "关键洞察"),
            ("3", "分析类型")  # 单变量、双变量、多变量
        ]
        
        x_positions = [1.5, 5.5, 9.5]
        for i, (number, label) in enumerate(metrics):
            # 数字
            num_box = slide.shapes.add_textbox(
                Inches(x_positions[i]), Inches(2), Inches(2.5), Inches(1.2)
            )
            num_frame = num_box.text_frame
            num_frame.text = number
            num_frame.paragraphs[0].font.size = Pt(60)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 标签
            label_box = slide.shapes.add_textbox(
                Inches(x_positions[i]), Inches(3.2), Inches(2.5), Inches(0.5)
            )
            label_frame = label_box.text_frame
            label_frame.text = label
            label_frame.paragraphs[0].font.size = Pt(18)
            label_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
            label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 核心结论
        conclusion_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(12.333), Inches(1.5)
        )
        conclusion_frame = conclusion_box.text_frame
        if result.insights:
            conclusion_frame.text = f"核心结论: {result.insights[0].get('description', '数据分析完成')[:80]}..."
        else:
            conclusion_frame.text = "核心结论: 基于多维度分析，数据质量良好，发现若干有价值的洞察。"
        conclusion_frame.paragraphs[0].font.size = Pt(20)
        conclusion_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
    
    def _add_data_overview_slide(self, result: AnalysisResult):
        """添加数据概览页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "数据概览"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        preprocessing = result.preprocessing
        structure = preprocessing.get('structure_understanding', {})
        columns = structure.get('columns', [])
        quality = preprocessing.get('quality_assessment', {})
        
        # 左侧：数据规模
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(5.5), Inches(3)
        )
        left_frame = left_box.text_frame
        left_frame.text = "数据规模"
        left_frame.paragraphs[0].font.size = Pt(24)
        left_frame.paragraphs[0].font.bold = True
        left_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
        
        p = left_frame.add_paragraph()
        p.text = f"• 总列数: {len(columns)}"
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        p = left_frame.add_paragraph()
        p.text = f"• 分析维度: {len(result.merged_dimensions)}"
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        # 右侧：数据质量
        right_box = slide.shapes.add_textbox(
            Inches(7), Inches(1.5), Inches(5.5), Inches(3)
        )
        right_frame = right_box.text_frame
        right_frame.text = "数据质量"
        right_frame.paragraphs[0].font.size = Pt(24)
        right_frame.paragraphs[0].font.bold = True
        right_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
        
        missing = quality.get('missing_values', {})
        p = right_frame.add_paragraph()
        p.text = f"• 缺失值: {len(missing)}列存在缺失"
        p.font.size = Pt(16)
        p.space_after = Pt(12)
        
        duplicates = quality.get('duplicates', {})
        p = right_frame.add_paragraph()
        p.text = f"• 重复数据: {duplicates.get('count', 0)}行"
        p.font.size = Pt(16)
        
        # 底部：分析维度
        dims_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.8), Inches(12.333), Inches(2)
        )
        dims_frame = dims_box.text_frame
        dims_frame.text = "分析维度来源"
        dims_frame.paragraphs[0].font.size = Pt(20)
        dims_frame.paragraphs[0].font.bold = True
        dims_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
        
        template_count = len([d for d in result.merged_dimensions if d.source == 'template'])
        ai_count = len([d for d in result.merged_dimensions if d.source == 'ai'])
        user_count = len([d for d in result.merged_dimensions if d.source == 'user'])
        
        p = dims_frame.add_paragraph()
        p.text = f"模板: {template_count}  |  AI自动: {ai_count}  |  用户: {user_count}"
        p.font.size = Pt(16)
    
    def _add_insight_slide(self, insight: Dict, index: int, chart_path: Optional[str] = None):
        """添加洞察页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"洞察 {index}: {insight.get('title', '')}"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # 左侧：详细描述
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.5)
        )
        left_frame = left_box.text_frame
        left_frame.text = insight.get('description', '')[:150] + "..."
        left_frame.paragraphs[0].font.size = Pt(14)
        left_frame.word_wrap = True
        
        # 关键发现
        findings_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(5.5), Inches(1.5)
        )
        findings_frame = findings_box.text_frame
        findings_frame.text = "关键发现"
        findings_frame.paragraphs[0].font.size = Pt(16)
        findings_frame.paragraphs[0].font.bold = True
        findings_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
        
        findings = insight.get('key_findings', [])[:2]
        for finding in findings:
            p = findings_frame.add_paragraph()
            p.text = f"• {finding[:50]}"
            p.font.size = Pt(12)
            p.level = 1
        
        # 右侧：图表或数据证据
        if chart_path and self._file_exists(chart_path):
            try:
                slide.shapes.add_picture(
                    chart_path,
                    Inches(6.5), Inches(1.3),
                    width=Inches(6)
                )
            except Exception:
                # 图表加载失败，显示数据证据
                self._add_data_evidence(slide, insight)
        else:
            # 显示数据证据
            self._add_data_evidence(slide, insight)
        
        # 底部：行动建议
        action_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.8)
        )
        action_frame = action_box.text_frame
        action_frame.text = f"建议行动: {insight.get('action', '继续监控')}"
        action_frame.paragraphs[0].font.size = Pt(16)
        action_frame.paragraphs[0].font.bold = True
        action_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)
        
        # 置信度和优先级
        meta_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.6), Inches(5), Inches(0.4)
        )
        meta_frame = meta_box.text_frame
        meta_frame.text = f"置信度: {insight.get('confidence', '中')}  |  优先级: {insight.get('priority', '中')}"
        meta_frame.paragraphs[0].font.size = Pt(12)
        meta_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    
    def _add_data_evidence(self, slide, insight: Dict):
        """添加数据证据（当没有图表时）"""
        evidence_box = slide.shapes.add_textbox(
            Inches(6.5), Inches(1.3), Inches(6), Inches(3)
        )
        evidence_frame = evidence_box.text_frame
        evidence_frame.text = "数据证据"
        evidence_frame.paragraphs[0].font.size = Pt(18)
        evidence_frame.paragraphs[0].font.bold = True
        evidence_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
        
        evidence = insight.get('data_evidence', [])
        for ev in evidence[:3]:
            p = evidence_frame.add_paragraph()
            p.text = f"• {ev}"
            p.font.size = Pt(14)
            p.level = 1
    
    def _add_conclusion_slide(self, result: AnalysisResult):
        """添加结论建议页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "结论与建议"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # 核心结论
        conclusion_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), Inches(12.333), Inches(1.5)
        )
        conclusion_frame = conclusion_box.text_frame
        conclusion_frame.text = "核心结论"
        conclusion_frame.paragraphs[0].font.size = Pt(22)
        conclusion_frame.paragraphs[0].font.bold = True
        conclusion_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
        
        p = conclusion_frame.add_paragraph()
        p.text = f"通过{len(result.merged_dimensions)}个分析维度，发现{len(result.insights)}个关键洞察"
        p.font.size = Pt(16)
        
        # 行动建议
        action_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(12.333), Inches(3)
        )
        action_frame = action_box.text_frame
        action_frame.text = "行动建议"
        action_frame.paragraphs[0].font.size = Pt(22)
        action_frame.paragraphs[0].font.bold = True
        action_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
        
        # 获取高优先级洞察的建议
        high_priority = [i for i in result.insights if i.get('priority') == '高'][:3]
        if not high_priority:
            high_priority = result.insights[:3]
        
        for i, insight in enumerate(high_priority, 1):
            p = action_frame.add_paragraph()
            p.text = f"{i}. {insight.get('title', '')}: {insight.get('action', '继续监控')}"
            p.font.size = Pt(16)
            p.space_after = Pt(12)
    
    def _add_thank_you_slide(self):
        """添加感谢页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 感谢文字
        thanks_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(2)
        )
        thanks_frame = thanks_box.text_frame
        thanks_frame.text = "感谢聆听"
        thanks_frame.paragraphs[0].font.size = Pt(54)
        thanks_frame.paragraphs[0].font.bold = True
        thanks_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        thanks_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(12.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "欢迎提问与讨论"
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _file_exists(self, path: str) -> bool:
        """检查文件是否存在"""
        import os
        return os.path.exists(path)
